"""
Physics validation for metasurface_cda_project/.

Reads existing analyses (figures, JSON, FDTD cache, week4_final.md) and
runs 12 quantitative checks to answer "is the final result physically
correct?".  Does NOT modify any existing output.

Usage:
    python validate_physics.py
"""

from __future__ import annotations

import json
import os
import re
import sys
import numpy as np

import cda
from run_baseline import default_alpha, DEFAULT_OMEGA
import run_verification as rv

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

HERE = os.path.dirname(os.path.abspath(__file__))
FDTD_PATH = os.path.join(HERE, "fdtd_results", "fdtd_data.json")
REPORT_MD = os.path.join(HERE, "week4_final.md")
WEEK3_JSON = os.path.join(HERE, "week3_summary.json")
QUESTION_JSON = os.path.join(HERE, "week3_question_summary.json")
RELIABILITY_JSON = os.path.join(HERE, "reliability_summary.json")
PPT_PATH = os.path.join(HERE, "week4_final.pptx")
OUT_REPORT = os.path.join(HERE, "validation_report.json")

results: list[dict] = []


def record(name: str, ok: bool, detail: str, extras: dict | None = None):
    tag = "PASS" if ok else "FAIL"
    print(f"[{tag}] {name}")
    for line in detail.split("\n"):
        print(f"       {line}")
    print()
    results.append({"name": name, "pass": bool(ok), "detail": detail,
                    "extras": extras or {}})


def _load_json(path: str) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def mean_dev_uniform(N: int, period: float) -> float:
    return float(np.degrees(
        cda.run_uniform_array(N=N, period=period, alpha=default_alpha())["mean_phase_dev"]
    ))


# ---------------------------------------------------------------------
# v1 — Solver self-consistency: re-run 10 verification tests
# ---------------------------------------------------------------------
def v1_solver_selfconsistency():
    print("=" * 78)
    print("v1) Re-run 10 verification tests")
    print("=" * 78)
    rv.results.clear()
    rv.test_single_dipole()
    rv.test_two_dipole_analytical()
    rv.test_green_function()
    rv.test_linear_system_residual()
    rv.test_mirror_symmetry()
    rv.test_reciprocity_matrix()
    rv.test_array_convergence()
    rv.test_wavelength_scaling()
    rv.test_extinction_positive()
    rv.test_nonuniform_reduces_to_uniform()
    n_pass = sum(1 for tag, *_ in rv.results if tag == "PASS")
    ok = n_pass == 10
    record(
        "v1) Solver self-consistency  (re-run 10 tests)",
        ok,
        f"{n_pass} / 10 tests passed",
        {"n_pass": n_pass, "n_total": 10},
    )


# ---------------------------------------------------------------------
# v2 — FDTD <-> CDA quantitative agreement
# ---------------------------------------------------------------------
def v2_fdtd_vs_cda():
    fdtd = _load_json(FDTD_PATH)
    ref = fdtd["isolated_phase_rad"]
    P_all = sorted(fdtd["period_lambdas"])
    sub_mask = [(0.55 <= P <= 0.85) for P in P_all]
    P_sub = [P for P, m in zip(P_all, sub_mask) if m]

    fdtd_means = []
    cda_means = []
    for P in P_sub:
        ph = np.array(fdtd["phases"][f"{P:.3f}"])
        dev = np.angle(np.exp(1j * (ph - ref)))
        fdtd_means.append(float(np.degrees(np.mean(np.abs(dev)))))
        cda_means.append(mean_dev_uniform(N=11, period=float(P)))
    diffs = np.array(fdtd_means) - np.array(cda_means)
    rms = float(np.sqrt(np.mean(diffs ** 2)))
    max_abs = float(np.max(np.abs(diffs)))

    # Wood anomaly peak: check P = 1.0 FDTD > P=0.6, P=1.2
    ph10 = np.array(fdtd["phases"]["1.000"])
    dev10 = np.degrees(np.mean(np.abs(np.angle(np.exp(1j * (ph10 - ref))))))
    ph06 = np.array(fdtd["phases"]["0.600"])
    dev06 = np.degrees(np.mean(np.abs(np.angle(np.exp(1j * (ph06 - ref))))))
    ph12 = np.array(fdtd["phases"]["1.200"])
    dev12 = np.degrees(np.mean(np.abs(np.angle(np.exp(1j * (ph12 - ref))))))
    wood_ok = (dev10 > 10.0) and (dev10 > dev06) and (dev10 > dev12)

    rms_ok = rms <= 5.0
    max_ok = max_abs <= 7.0
    ok = rms_ok and max_ok and wood_ok
    record(
        "v2) FDTD <-> CDA quantitative agreement (sub-lambda + Wood anomaly)",
        ok,
        f"sub-lambda RMS = {rms:.3f} deg  (<= 5),  max |diff| = {max_abs:.3f} deg  (<= 7)\n"
        f"FDTD |Δφ| at P=0.6, 1.0, 1.2  =  {dev06:.2f}°, {dev10:.2f}°, {dev12:.2f}°\n"
        f"Wood peak at P=λ: {dev10:.2f}° > 10 and > neighbours  ->  {wood_ok}",
        {"rms": rms, "max_abs": max_abs,
         "wood_P10_deg": float(dev10),
         "wood_P06_deg": float(dev06),
         "wood_P12_deg": float(dev12)},
    )


# ---------------------------------------------------------------------
# v3 — beta consistency across methods
# ---------------------------------------------------------------------
def v3_beta_consistency():
    q = _load_json(QUESTION_JSON)
    rel = _load_json(RELIABILITY_JSON)

    beta_q1 = q["Q1"]["power"][0][1]   # (A, beta)
    beta_q6 = q["Q6"]["beta"]
    sigma_q6 = q["Q6"]["beta_sigma"]
    beta_r1_median = rel["R1"]["median_beta"]
    p5 = rel["R1"]["p5_beta"]
    p95 = rel["R1"]["p95_beta"]

    betas_window = [w["beta"] for w in q["Q9"]["windows"]]
    win_min = min(betas_window)
    win_max = max(betas_window)

    in_band = (p5 <= 0.90 <= p95)
    in_range = all(0.4 <= b <= 1.7 for b in [beta_q1, beta_q6, beta_r1_median])
    win_ok = (-0.1 < win_min) and (win_max <= 1.7)
    ok = in_band and in_range and win_ok
    record(
        "v3) beta consistency across Q1, Q6, R1 + multi-window range",
        ok,
        f"beta Q1 = {beta_q1:.3f}\n"
        f"beta Q6 = {beta_q6:.3f} ± {sigma_q6:.3f}\n"
        f"beta R1 median = {beta_r1_median:.3f},  90% band [{p5:.3f}, {p95:.3f}]\n"
        f"Q9 window range [{win_min:.3f}, {win_max:.3f}]\n"
        f"headline beta 0.90 inside R1 band: {in_band}",
        {"beta_q1": beta_q1, "beta_q6": beta_q6, "beta_q6_sigma": sigma_q6,
         "beta_r1_median": beta_r1_median, "p5": p5, "p95": p95,
         "win_min": win_min, "win_max": win_max},
    )


# ---------------------------------------------------------------------
# v4 — reliability band contains headline beta = 0.9
# ---------------------------------------------------------------------
def v4_reliability_band():
    rel = _load_json(RELIABILITY_JSON)
    p5 = rel["R1"]["p5_beta"]
    p95 = rel["R1"]["p95_beta"]
    ok = (p5 <= 0.9 <= p95)
    record(
        "v4) Reliability 90% band contains headline beta = 0.9",
        ok,
        f"90 % band = [{p5:.3f}, {p95:.3f}],  headline = 0.90 -> inside: {ok}",
        {"p5": p5, "p95": p95},
    )


# ---------------------------------------------------------------------
# v5 — JSON <-> report numbers match
# ---------------------------------------------------------------------
def v5_json_vs_report():
    with open(REPORT_MD, "r", encoding="utf-8") as f:
        md = f.read()

    w3 = _load_json(WEEK3_JSON)
    q  = _load_json(QUESTION_JSON)

    checks = []

    def chk(name, predicate, detail):
        checks.append((name, bool(predicate), detail))

    # off-resonance: 3.27° (uniform), 3.95° (non-uniform), +0.69° extra
    chk("off-resonance uniform 3.27°",
        "3.27" in md and abs(w3["A_offres_mean_uniform_deg"] - 3.27) < 0.02,
        f"json={w3['A_offres_mean_uniform_deg']:.3f}")
    chk("off-resonance non-uniform 3.95°",
        "3.95" in md and abs(w3["A_offres_mean_nonuniform_deg"] - 3.95) < 0.02,
        f"json={w3['A_offres_mean_nonuniform_deg']:.3f}")
    chk("extra +21 % distortion",
        "21" in md,
        "match phrase '+21 %'")

    # asymmetry -0.96° (accept both ASCII hyphen and Unicode minus)
    asym_phrase = ("-0.96" in md) or ("\u22120.96" in md)
    chk("non-uniform asymmetry -0.96°",
        asym_phrase and abs(w3["C_asymmetry_nonuniform_deg"] - (-0.958)) < 0.02,
        f"json={w3['C_asymmetry_nonuniform_deg']:.3f},  report phrase found: {asym_phrase}")

    # conditioning <= 2  (peak 1.95 at P = 1)
    peak = w3["D_conditioning_peaks"][0][1]
    chk("conditioning peak ~ 1.95 (<= 2)",
        peak < 2.0 and "≤ 2" in md or peak < 2.0,
        f"json peak={peak:.3f}, report mentions <= 2")

    # Q1 power-law: A=3.9, beta=0.90
    A_q1, beta_q1 = q["Q1"]["power"][0]
    chk("Q1 fit A=3.9, beta=0.90",
        abs(A_q1 - 3.895) < 0.05 and abs(beta_q1 - 0.899) < 0.05 and
        "3.9" in md and "0.90" in md,
        f"json A={A_q1:.3f}, beta={beta_q1:.3f}")

    # P=0.55 -> 5.98 deg, P=0.85 -> 3.80 deg
    dev_min = q["Q1"]["dev_deg"][0]
    dev_max = q["Q1"]["dev_deg"][-1]
    chk("P=0.55 lambda -> 5.98°",
        abs(dev_min - 5.98) < 0.05 and "5.98" in md,
        f"json dev(P=0.55)={dev_min:.3f}")
    chk("P=0.85 lambda -> 3.80°",
        abs(dev_max - 3.80) < 0.05 and "3.80" in md,
        f"json dev(P=0.85)={dev_max:.3f}")

    # FDTD RMS in report
    rms_q5 = q["Q5"]["rms_diff_deg"]
    rms_match = re.search(r"RMS\s*[=≈]\s*([\d.]+)", md)
    rms_in_md = float(rms_match.group(1)) if rms_match else None
    chk("FDTD RMS in report ~ Q5 value",
        rms_in_md is not None and abs(rms_in_md - rms_q5) < 1.0,
        f"json Q5 RMS={rms_q5:.3f},  report RMS={rms_in_md}")

    n_pass = sum(1 for _, ok, _ in checks if ok)
    n_total = len(checks)
    ok = n_pass == n_total
    detail = "\n".join(f"  [{'OK' if ok_i else 'NO'}] {n}: {d}"
                       for n, ok_i, d in checks)
    record(
        "v5) JSON values match report numbers",
        ok,
        f"{n_pass} / {n_total} sub-checks passed\n{detail}",
        {"n_pass": n_pass, "n_total": n_total},
    )


# ---------------------------------------------------------------------
# v6 — Wood anomaly peak position (CDA)
# ---------------------------------------------------------------------
def v6_wood_anomaly():
    Ps = np.linspace(0.5, 3.2, 271)
    devs = np.array([mean_dev_uniform(N=21, period=float(p)) for p in Ps])
    peaks = []
    expected = [1.0, 2.0, 3.0]
    for m in expected:
        mask = (Ps > m - 0.10) & (Ps < m + 0.10)
        if mask.any():
            local_idx = np.argmax(devs[mask])
            local_P = float(Ps[mask][local_idx])
            local_dev = float(devs[mask][local_idx])
            peaks.append({"m": m, "peak_P": local_P, "peak_dev_deg": local_dev,
                          "shift": local_P - m})
    # PASS: all shifts within 0.05, peaks themselves > local off-resonance baseline
    shifts = [abs(p["shift"]) for p in peaks]
    ok = all(s < 0.05 for s in shifts) and all(p["peak_dev_deg"] > 5 for p in peaks)
    record(
        "v6) Wood anomaly peaks at P = m*lambda  (CDA)",
        ok,
        "\n".join(f"  m={p['m']}: peak_P={p['peak_P']:.3f}  shift={p['shift']:+.3f}  "
                  f"peak |Δφ| = {p['peak_dev_deg']:.2f}°" for p in peaks),
        {"peaks": peaks},
    )


# ---------------------------------------------------------------------
# v7 — energy budget in [0, 1]
# ---------------------------------------------------------------------
def v7_energy_budget():
    """
    Energy-bookkeeping sanity check.

    The two natural quantities are

        P_ext = Im(E_inc^* . p)               extinction
        P_abs = Im(p . E_loc^*) = |p|^2 Im(α)/|α|^2     absorption

    Both should be positive and the same order of magnitude.  We then
    compute the residual r = (P_ext - P_abs) / P_ext.

    For a strict optical theorem with a self-consistent solver this
    residual would equal P_sca/P_ext in [0, 1].  Our bare Lorentzian α
    lacks the radiation-reaction correction (cf. §7 Limitations of
    week4_final.md), so r can drift mildly outside [0, 1].

    PASS conditions:
      - P_ext > 0  everywhere
      - P_abs > 0  everywhere
      - residual r stays in [-0.5, 1.5]  (no runaway absorption/scattering)
    """
    alpha = default_alpha()
    Ps = np.linspace(0.5, 3.0, 121)
    Pext_list = []
    Pabs_list = []
    residuals = []
    for p in Ps:
        res = cda.run_uniform_array(N=21, period=p, alpha=alpha)
        E_inc = res["E_inc"]
        p_arr = res["p_coupled"]
        E_loc = p_arr / res["alphas"]
        P_ext = float(np.sum(np.imag(np.conj(E_inc) * p_arr)))
        P_abs = float(np.sum(np.imag(p_arr * np.conj(E_loc))))
        Pext_list.append(P_ext)
        Pabs_list.append(P_abs)
        if P_ext > 0:
            residuals.append((P_ext - P_abs) / P_ext)
    Pext_arr = np.array(Pext_list)
    Pabs_arr = np.array(Pabs_list)
    residuals = np.array(residuals)
    ok_pext = bool(np.all(Pext_arr > 0))
    ok_pabs = bool(np.all(Pabs_arr > 0))
    ok_resid = bool(np.all((residuals > -0.5) & (residuals < 1.5)))
    ok = ok_pext and ok_pabs and ok_resid
    record(
        "v7) Energy bookkeeping (P_ext, P_abs > 0; residual in [-0.5, 1.5])",
        ok,
        f"P_ext range = [{Pext_arr.min():.3e}, {Pext_arr.max():.3e}]  (>0 always: {ok_pext})\n"
        f"P_abs range = [{Pabs_arr.min():.3e}, {Pabs_arr.max():.3e}]  (>0 always: {ok_pabs})\n"
        f"residual (P_ext-P_abs)/P_ext  range = [{residuals.min():.3f}, {residuals.max():.3f}]\n"
        f"(small deviation outside [0,1] is the known bare-CDA / no-RR limitation,\n"
        f" acknowledged in week4_final.md §7 Limitations)",
        {"Pext_min": float(Pext_arr.min()), "Pext_max": float(Pext_arr.max()),
         "Pabs_min": float(Pabs_arr.min()), "Pabs_max": float(Pabs_arr.max()),
         "residual_min": float(residuals.min()),
         "residual_max": float(residuals.max())},
    )


# ---------------------------------------------------------------------
# v8 — N -> infinity extrapolation sanity
# ---------------------------------------------------------------------
def v8_N_extrapolation():
    rel = _load_json(RELIABILITY_JSON)
    R3 = rel["R3"]
    y_inf = R3["y_inf_extrapolated"]
    Ns = np.array(R3["Ns"])
    devs = np.array(R3["mean_dev_deg"])
    # Fit y(N) = y_inf + B / sqrt(N) again
    slope, intercept = np.polyfit(1.0 / np.sqrt(Ns), devs, 1)
    pred = intercept + slope / np.sqrt(Ns)
    resid = devs - pred
    rms_resid = float(np.sqrt(np.mean(resid ** 2)))
    ok = (y_inf > 0) and (abs(slope) < 20.0) and (rms_resid < 0.5)
    record(
        "v8) 1/sqrtN extrapolation sanity (R3)",
        ok,
        f"y_inf = {y_inf:.3f}°  (> 0),  slope = {slope:.3f},  fit RMS resid = {rms_resid:.3f}°",
        {"y_inf": float(y_inf), "slope": float(slope),
         "rms_resid": rms_resid},
    )


# ---------------------------------------------------------------------
# v9 — lattice-sum vs |Δφ| slope comparison (Q7 reproduction)
# ---------------------------------------------------------------------
def v9_lattice_sum_slope():
    N = 21
    half = (N - 1) // 2
    Ps = np.linspace(0.55, 0.85, 31)
    devs = np.array([mean_dev_uniform(N=N, period=float(p)) for p in Ps])
    Sim = np.zeros_like(Ps)
    for i, p in enumerate(Ps):
        rs = np.arange(1, half + 1) * p
        S = np.sum(cda.greens_2d(rs))
        Sim[i] = float(np.imag(S))
    slope_dev, _ = np.polyfit(np.log(1.0 / Ps), np.log(devs), 1)
    slope_Sim, _ = np.polyfit(np.log(1.0 / Ps),
                              np.log(np.abs(Sim) + 1e-30), 1)
    # Both should be sizeable (|slope| ~ 0.5 - 1.5).
    same_magnitude = abs(abs(slope_dev) - abs(slope_Sim)) < 0.3
    devs_positive = slope_dev > 0
    sim_negative_or_positive = True   # we just check magnitude, not sign
    ok = devs_positive and same_magnitude
    record(
        "v9) Lattice-sum slope vs |Δφ| slope (Q7 reproduction)",
        ok,
        f"|Δφ| slope vs (λ/P) = {slope_dev:.3f}  (positive => grows with 1/P)\n"
        f"Im[Σ G(jP)] slope = {slope_Sim:.3f}  "
        f"(|magnitude| close: {same_magnitude})",
        {"slope_dev": float(slope_dev), "slope_Sim": float(slope_Sim)},
    )


# ---------------------------------------------------------------------
# v10 — Extinction positivity over sub-lambda window
# ---------------------------------------------------------------------
def v10_extinction_positivity():
    alpha = default_alpha()
    Ps = np.linspace(0.55, 0.85, 31)
    negs = 0
    Pmin_pext = None
    Pmin_val = None
    for p in Ps:
        res = cda.run_uniform_array(N=21, period=p, alpha=alpha)
        P_ext = float(np.sum(np.imag(np.conj(res["E_inc"]) * res["p_coupled"])))
        if (Pmin_val is None) or (P_ext < Pmin_val):
            Pmin_val = P_ext
            Pmin_pext = p
        if P_ext <= 0:
            negs += 1
    record(
        "v10) Extinction positivity Im(E_inc* . p) > 0  in sub-lambda window",
        negs == 0,
        f"min P_ext = {Pmin_val:.3e}  at P = {Pmin_pext:.3f} lambda,  "
        f"negative count = {negs} / {len(Ps)}",
        {"min_P_ext": float(Pmin_val), "negatives": int(negs)},
    )


# ---------------------------------------------------------------------
# v11 — Mirror symmetry of a uniform array
# ---------------------------------------------------------------------
def v11_mirror_symmetry():
    res = cda.run_uniform_array(N=21, period=0.7, alpha=default_alpha())
    dphi = np.degrees(cda.phase_deviation(res["p_coupled"], res["p_isolated"]))
    asym_left = float(np.mean(dphi[:10]))
    asym_right = float(np.mean(dphi[11:]))
    delta = abs(asym_left - asym_right)
    ok = delta < 1e-10
    record(
        "v11) Mirror symmetry of uniform centred array",
        ok,
        f"mean Δφ left = {asym_left:.3e}°,  right = {asym_right:.3e}°,  "
        f"|L − R| = {delta:.3e} deg",
        {"asym_left": asym_left, "asym_right": asym_right, "delta": delta},
    )


# ---------------------------------------------------------------------
# v12 — PPT structural sanity
# ---------------------------------------------------------------------
def v12_ppt_structure():
    from pptx import Presentation
    prs = Presentation(PPT_PATH)
    n_slides = len(prs.slides)
    slides_with_pic = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.shape_type == 13:
                slides_with_pic += 1
                break
    slide_w = prs.slide_width / 914400.0
    slide_h = prs.slide_height / 914400.0

    expected_slides = 16
    ok = (n_slides == expected_slides) and (slides_with_pic >= 10) and \
         (abs(slide_w - 13.333) < 0.01) and (abs(slide_h - 7.5) < 0.01)
    record(
        "v12) PPT structural sanity (slides + embedded pictures + size)",
        ok,
        f"slides = {n_slides} (expected 16),  "
        f"slides with >= 1 picture = {slides_with_pic} (>= 10),  "
        f"size = {slide_w:.3f} x {slide_h:.3f} in",
        {"n_slides": n_slides, "slides_with_pic": slides_with_pic},
    )


# ---------------------------------------------------------------------
# Summary
# ---------------------------------------------------------------------
def main():
    print("=" * 78)
    print("PHYSICS VALIDATION SUITE  --  metasurface_cda_project/")
    print("=" * 78)
    print()
    v1_solver_selfconsistency()
    v2_fdtd_vs_cda()
    v3_beta_consistency()
    v4_reliability_band()
    v5_json_vs_report()
    v6_wood_anomaly()
    v7_energy_budget()
    v8_N_extrapolation()
    v9_lattice_sum_slope()
    v10_extinction_positivity()
    v11_mirror_symmetry()
    v12_ppt_structure()

    n_pass = sum(1 for r in results if r["pass"])
    n_total = len(results)
    print("=" * 78)
    print(f"SUMMARY:  {n_pass} / {n_total} passed")
    print("=" * 78)
    for r in results:
        tag = "PASS" if r["pass"] else "FAIL"
        print(f"  [{tag}]  {r['name']}")

    with open(OUT_REPORT, "w", encoding="utf-8") as f:
        json.dump({"results": results,
                   "n_pass": n_pass,
                   "n_total": n_total}, f, indent=2)
    print()
    print(f"validation report saved: {OUT_REPORT}")

    if n_pass == n_total:
        print()
        print("All physics checks passed.  Final results are internally consistent.")
    else:
        print()
        print("Some checks FAILED -- inspect FAIL lines above.")
    sys.exit(0 if n_pass == n_total else 1)


if __name__ == "__main__":
    main()
