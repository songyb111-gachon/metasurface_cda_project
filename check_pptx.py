"""
Validate the generated PPTX file.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Emu

PPTX_NAME = sys.argv[1] if len(sys.argv) > 1 else "week4_final.pptx"
HERE = os.path.dirname(os.path.abspath(__file__))
PPTX_PATH = os.path.join(HERE, PPTX_NAME)
FIG_DIR = os.path.join(HERE, "figures")

issues: list[str] = []
warnings: list[str] = []


def emu_to_inches(v: int) -> float:
    return v / 914400.0


def check_file_exists():
    if not os.path.exists(PPTX_PATH):
        issues.append(f"PPTX file does not exist: {PPTX_PATH}")
        return False
    size = os.path.getsize(PPTX_PATH)
    print(f"[OK]  File exists: {PPTX_PATH}  ({size/1024:.1f} KB)")
    return True


def main():
    if not check_file_exists():
        return

    prs = Presentation(PPTX_PATH)
    n_slides = len(prs.slides)
    slide_w = emu_to_inches(prs.slide_width)
    slide_h = emu_to_inches(prs.slide_height)

    print(f"[OK]  Slide count: {n_slides}")
    print(f"[OK]  Slide size : {slide_w:.2f} x {slide_h:.2f} in")
    print()
    print("=" * 90)
    print(f"{'Slide':>5} | {'Shapes':>6} | {'Text':>5} | {'Pics':>5} | "
          f"{'Empty':>5} | {'Overflow':>8} | First text")
    print("=" * 90)

    for idx, slide in enumerate(prs.slides):
        n_shapes = len(slide.shapes)
        n_text = n_pics = n_empty = n_overflow = 0
        first_text = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                tf = sh.text_frame
                full = tf.text.strip()
                if not full:
                    n_empty += 1
                else:
                    n_text += 1
                    if not first_text:
                        snippet = full[:42].replace("\n", " ")
                        first_text = snippet.encode("ascii", "replace").decode("ascii")
            if sh.shape_type == 13:
                n_pics += 1
            try:
                left = emu_to_inches(sh.left or 0)
                top = emu_to_inches(sh.top or 0)
                w = emu_to_inches(sh.width or 0)
                h = emu_to_inches(sh.height or 0)
                if (left + w) > slide_w + 0.05 or (top + h) > slide_h + 0.05:
                    n_overflow += 1
            except (AttributeError, TypeError):
                pass
        print(f"{idx:>5} | {n_shapes:>6} | {n_text:>5} | {n_pics:>5} | "
              f"{n_empty:>5} | {n_overflow:>8} | {first_text}")
        if n_overflow:
            warnings.append(
                f"Slide {idx}: {n_overflow} shape(s) overflow slide bounds")

    print()
    print("=" * 90)
    print(f"SUMMARY:  {len(issues)} issues,  {len(warnings)} warnings")
    print("=" * 90)
    for x in issues:
        print(f"  [ISSUE]    {x}")
    for x in warnings:
        print(f"  [WARNING]  {x}")
    if not issues and not warnings:
        print("  [OK]  No issues found.")


if __name__ == "__main__":
    main()
